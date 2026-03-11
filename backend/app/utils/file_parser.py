import base64
import io
import logging
import os

logger = logging.getLogger(__name__)

# Candidate Tesseract executable paths — checked in order when PATH lookup fails.
_TESSERACT_CANDIDATES = [
    # Linux (Render / Ubuntu)
    "/usr/bin/tesseract",
    "/usr/local/bin/tesseract",
    # Windows (local dev)
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
    r"C:\Tesseract-OCR\tesseract.exe",
]

# Groq vision models tried in order — first available wins.
_GROQ_VISION_MODELS = [
    "meta-llama/llama-4-scout-17b-16e-instruct",   # Llama 4 Scout — best, current
    "meta-llama/llama-4-maverick-17b-128e-instruct", # Llama 4 Maverick — fallback
    "llama-3.2-11b-vision-preview",                  # Llama 3.2 — legacy fallback
]


def _find_tesseract_cmd() -> str | None:
    """Return the first candidate tesseract.exe path that exists, or None."""
    for candidate in _TESSERACT_CANDIDATES:
        if os.path.isfile(candidate):
            return candidate
    return None


_OCR_PAGE_LIMIT = 20  # max pages sent for OCR — covers most academic documents


def _render_pdf_pages(file_bytes: bytes, max_pages: int = _OCR_PAGE_LIMIT) -> list:
    """
    Render up to max_pages PDF pages as JPEG images using pymupdf.
    Uses 3× zoom (216 DPI) for better quality on scanned/handwritten documents.
    """
    import fitz  # pymupdf

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    mat = fitz.Matrix(3, 3)  # 3× zoom ≈ 216 DPI — sharper for handwritten/scanned docs
    pages_to_render = min(len(doc), max_pages)
    images = []
    for page_num in range(pages_to_render):
        pix = doc[page_num].get_pixmap(matrix=mat)
        images.append(pix.tobytes("jpeg", jpg_quality=90))
        del pix  # free pixmap memory immediately
    doc.close()
    return images


def _ocr_with_tesseract(page_images: list) -> str:
    """
    Try pytesseract OCR on pre-rendered page images.
    Returns '' if Tesseract is not installed or pytesseract is missing.
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""

    tesseract_cmd = _find_tesseract_cmd()
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

    parts = []
    for i, img_bytes in enumerate(page_images):
        try:
            img = Image.open(io.BytesIO(img_bytes))
            text = pytesseract.image_to_string(img)
            img.close()
            if text.strip():
                parts.append(text.strip())
        except Exception as exc:
            logger.warning("Tesseract OCR failed on page %d: %s", i + 1, exc)

    return "\n\n".join(parts)


def _ocr_with_groq(page_images: list) -> str:
    """
    Cloud OCR fallback: send page images to a Groq vision model one at a time.
    Processes pages sequentially and releases each image from memory after use.
    """
    api_key = os.environ.get("GROQ_API_KEY", "")
    if not api_key:
        logger.warning("GROQ_API_KEY not set — Groq vision OCR unavailable")
        return ""

    try:
        from groq import Groq
        client = Groq(api_key=api_key)
    except Exception as exc:
        logger.warning("Groq client init failed: %s", exc)
        return ""

    def _ocr_page(model: str, img_bytes: bytes, page_num: int) -> str:
        img_b64 = base64.b64encode(img_bytes).decode()
        response = client.chat.completions.create(
            model=model,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"},
                    },
                    {
                        "type": "text",
                        "text": (
                            "This is a page from an academic document — it may be typed, printed, "
                            "or handwritten lecture notes. Extract ALL text you can read from this page. "
                            "For handwritten content, do your best to interpret the writing accurately. "
                            "Preserve headings, numbered lists, equations, and structure where possible. "
                            "IMPORTANT: If any text appears visually emphasized — bold, highlighted, "
                            "underlined, circled, starred, or otherwise marked — wrap it in **...** "
                            "(e.g. **Option C**). This is critical for detecting correct answers. "
                            "Output only the extracted text — no commentary, no descriptions of the image."
                        ),
                    },
                ],
            }],
            max_tokens=2048,
        )
        del img_b64  # free the large base64 string immediately
        return response.choices[0].message.content.strip()

    for model in _GROQ_VISION_MODELS:
        parts = []
        failed = False
        for i, img_bytes in enumerate(page_images):
            try:
                text = _ocr_page(model, img_bytes, i + 1)
                if text:
                    parts.append(text)
            except Exception as exc:
                logger.warning("Groq vision OCR failed on page %d with %s: %s", i + 1, model, exc)
                failed = True
                break
        if parts and not failed:
            logger.info("Groq vision OCR succeeded with model %s", model)
            return "\n\n".join(parts)

    return ""


def _ocr_pdf(file_bytes: bytes) -> str:
    """
    OCR fallback for scanned / image-only PDFs.

    Pipeline:
    1. Render pages to images via pymupdf  (no Poppler system package needed)
    2. Try pytesseract                     (fast, local — requires Tesseract binary)
    3. Fall back to Groq vision API        (cloud — zero system dependencies)

    Returns empty string if all strategies fail.
    """
    try:
        page_images = _render_pdf_pages(file_bytes)
    except Exception as exc:
        logger.warning("PDF page rendering failed: %s", exc)
        return ""

    if not page_images:
        return ""

    # Strategy 1: local Tesseract
    text = _ocr_with_tesseract(page_images)
    if text:
        logger.info("Tesseract OCR succeeded — %d chars extracted.", len(text))
        return text

    # Strategy 2: Groq vision API
    logger.info("Tesseract unavailable — trying Groq vision OCR.")
    text = _ocr_with_groq(page_images)
    if text:
        logger.info("Groq vision OCR succeeded — %d chars extracted.", len(text))
    return text


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Extract text from a PDF.
    Strategy 1: pdfplumber  (fast, works for PDFs with a text layer)
    Strategy 2: OCR         (for scanned / image-only PDFs — pymupdf + Tesseract or Groq)
    """
    # --- Strategy 1: pdfplumber ---
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        text = "\n\n".join(text_parts).strip()
    except Exception as e:
        raise ValueError(f"Failed to open PDF: {str(e)}")

    if text:
        return text

    # --- Strategy 2: OCR fallback ---
    logger.info("pdfplumber returned no text — attempting OCR fallback.")

    # Check total page count so we can warn when truncation occurs
    try:
        import fitz
        total_pages = fitz.open(stream=file_bytes, filetype="pdf").page_count
    except Exception:
        total_pages = 0

    text = _ocr_pdf(file_bytes).strip()

    if text:
        logger.info("OCR succeeded — extracted %d characters.", len(text))
        if total_pages > _OCR_PAGE_LIMIT:
            text += (
                f"\n\n[Note: This document has {total_pages} pages. "
                f"OCR was applied to the first {_OCR_PAGE_LIMIT} pages only.]"
            )
        return text

    raise ValueError(
        "Could not extract text from this PDF. "
        "If it is a scanned document, please ensure your internet connection is active "
        "so the AI-powered OCR can process it, then try again."
    )


def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extract text from a DOCX file.
    Bold runs are wrapped in **...** so the AI can detect visually marked answers.
    """
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []
        for p in doc.paragraphs:
            if not p.text.strip():
                continue
            # Rebuild paragraph run-by-run, wrapping bold runs in **...**
            parts = []
            for run in p.runs:
                if not run.text:
                    continue
                if run.bold:
                    parts.append(f"**{run.text}**")
                else:
                    parts.append(run.text)
            line = "".join(parts).strip()
            if line:
                paragraphs.append(line)
        text = "\n\n".join(paragraphs).strip()
        if not text:
            raise ValueError(
                "This DOCX file contains no extractable text. "
                "Please ensure the document has readable text content."
            )
        return text
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Failed to parse DOCX: {str(e)}")


def render_pptx_visual_html(file_bytes: bytes) -> str:
    """
    Convert a PPTX into visual HTML slides preserving layout, positions,
    background colours, font sizes and font colours from the original file.
    Each slide is rendered as a positioned container using cqw units so text
    scales correctly when the panel is resized.
    """
    import html as _html
    try:
        from pptx import Presentation
        from pptx.enum.dml import PP_ALIGN
    except ImportError:
        return ""

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception:
        return ""

    sw = prs.slide_width  or 9144000   # EMUs
    sh = prs.slide_height or 6858000   # EMUs
    # aspect ratio expressed as padding-bottom % so the container is always 16:9 (or 4:3)
    ar_pct = round((sh / sw) * 100, 3)

    # 1 pt at 96 dpi on a 10"-wide slide = (pt / (sw/914400 * 96)) * 100 cqw
    # Simplified: cqw_per_pt = 914400 * 100 / (sw * 96 / 72)
    # At 9144000 EMU wide: 914400*100 / (9144000*96/72) ≈ 0.1042 cqw/pt
    emu_to_pct_w = 100.0 / sw
    emu_to_pct_h = 100.0 / sh
    pt_to_cqw    = (914400 * 100.0) / (sw * (96.0 / 72.0))

    def _rgb(color_obj):
        try:
            return f"#{color_obj.rgb}"
        except Exception:
            return None

    def _slide_bg(slide):
        """Try to get background colour from slide → layout → master."""
        for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
            try:
                fill = src.background.fill
                if fill.type is not None:
                    c = _rgb(fill.fore_color)
                    if c:
                        return c
            except Exception:
                pass
        return "#ffffff"

    slides_html = []
    for s_idx, slide in enumerate(prs.slides):
        bg = _slide_bg(slide)
        shapes_html = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                lp = shape.left  / sw * 100
                tp = shape.top   / sh * 100
                wp = shape.width / sw * 100
                hp = shape.height/ sh * 100
            except Exception:
                continue

            paras = []
            for para in shape.text_frame.paragraphs:
                runs_text = "".join(r.text for r in para.runs).strip()
                if not runs_text:
                    continue

                # Font properties — read from first run, fall back to para/shape defaults
                fsize_cqw = round(18 * pt_to_cqw, 2)   # 18 pt default
                fcolor    = "#000000"
                fbold     = False
                fitalic   = False
                falign    = "left"

                try:
                    if para.alignment == PP_ALIGN.CENTER:
                        falign = "center"
                    elif para.alignment == PP_ALIGN.RIGHT:
                        falign = "right"
                except Exception:
                    pass

                for run in para.runs:
                    try:
                        if run.font.size:
                            fsize_cqw = round(run.font.size.pt * pt_to_cqw, 2)
                    except Exception:
                        pass
                    try:
                        if run.font.bold:
                            fbold = True
                    except Exception:
                        pass
                    try:
                        if run.font.italic:
                            fitalic = True
                    except Exception:
                        pass
                    try:
                        c = _rgb(run.font.color)
                        if c:
                            fcolor = c
                    except Exception:
                        pass
                    break  # properties from first run apply to whole para

                style = (
                    f"margin:0 0 0.4em;"
                    f"font-size:{fsize_cqw}cqw;"
                    f"color:{fcolor};"
                    f"font-weight:{'700' if fbold else '400'};"
                    f"font-style:{'italic' if fitalic else 'normal'};"
                    f"text-align:{falign};"
                    f"line-height:1.25;word-break:break-word;"
                )
                paras.append(f'<p style="{style}">{_html.escape(runs_text)}</p>')

            if paras:
                shapes_html.append(
                    f'<div style="position:absolute;left:{lp:.2f}%;top:{tp:.2f}%;'
                    f'width:{wp:.2f}%;height:{hp:.2f}%;overflow:hidden;padding:0.4%;">'
                    + "".join(paras)
                    + "</div>"
                )

        slide_num = s_idx + 1
        slides_html.append(
            f'<div class="pptx-vs-outer">'
            f'<div class="pptx-vs-sizer" style="padding-bottom:{ar_pct}%;">'
            f'<div class="pptx-vs-inner" style="background:{bg};">'
            f'<span class="pptx-vs-num">SLIDE {slide_num}</span>'
            + "".join(shapes_html)
            + "</div></div></div>"
        )

    return "\n".join(slides_html)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text content from a PPTX (PowerPoint) file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide in prs.slides:
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text.strip())
                        if line.strip():
                            slide_parts.append(line.strip())
            if slide_parts:
                parts.append("\n".join(slide_parts))
        text = "\n\n".join(parts).strip()
        if not text:
            raise ValueError(
                "This PowerPoint file contains no extractable text. "
                "Please ensure the presentation has readable text content."
            )
        return text
    except ValueError:
        raise
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {str(e)}")


def extract_text_from_image(file_bytes: bytes) -> str:
    """
    Extract text from an image file (PNG, JPG, JPEG, WEBP, GIF) via OCR.

    Pipeline:
    1. Convert image to JPEG bytes using Pillow (normalise format for OCR engines).
    2. Try Tesseract (local, fast).
    3. Fall back to Groq vision API (cloud).

    Returns empty string if no text is found or OCR is unavailable.
    """
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(file_bytes))
        buf = io.BytesIO()
        img.convert("RGB").save(buf, "JPEG", quality=90)
        page_images = [buf.getvalue()]
    except Exception as exc:
        logger.warning("Image open/convert failed: %s", exc)
        return ""

    text = _ocr_with_tesseract(page_images)
    if text:
        logger.info("Tesseract OCR succeeded on image — %d chars.", len(text))
        return text

    logger.info("Tesseract unavailable for image — trying Groq vision OCR.")
    text = _ocr_with_groq(page_images)
    if text:
        logger.info("Groq vision OCR succeeded on image — %d chars.", len(text))
    return text


def truncate_text(text: str, max_chars: int = 8000) -> str:
    """Truncate text to avoid exceeding LLM context limits."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n... [content truncated]"
