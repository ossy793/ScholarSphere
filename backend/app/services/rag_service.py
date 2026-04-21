"""
RAG (Retrieval-Augmented Generation) utilities.

Flow:
  1. Extract text page-by-page from an uploaded document
  2. Split pages into overlapping text chunks, each carrying page metadata
  3. Embed chunks with fastembed (local, no API key required)
  4. Retrieve the top-k most relevant chunks for a query using cosine similarity
"""
import io
import logging
import re
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

_CHUNK_SIZE    = 900   # target characters per chunk
_CHUNK_OVERLAP = 120   # overlap between adjacent chunks so context isn't lost at boundaries


# ── Chapter / heading detector ────────────────────────────────────────────────

def _detect_heading(text: str) -> Optional[str]:
    """
    Return the first line if it looks like a chapter/section heading, else None.
    Matched patterns:
      - "Chapter 3 – Photosynthesis"
      - "UNIT 2: THERMODYNAMICS"
      - "1. Introduction"
      - Any ALL-CAPS line < 80 chars
    """
    first_line = text.strip().split("\n")[0].strip()
    if not first_line or len(first_line) > 100:
        return None
    if re.match(r"^(chapter|section|unit|part|module|topic)\s", first_line, re.IGNORECASE):
        return first_line[:80]
    if re.match(r"^\d+[\.\)]\s+\S.{2,60}$", first_line):
        return first_line[:80]
    if first_line == first_line.upper() and len(first_line) >= 4:
        return first_line[:80]
    return None


# ── Page-level text extraction ────────────────────────────────────────────────

def extract_pages_from_pdf(file_bytes: bytes) -> List[Dict]:
    """
    Extract text one page at a time from a PDF.
    Returns list of {"page": int (1-indexed), "text": str}.
    Falls back to empty list on failure (caller should handle gracefully).
    """
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = (page.extract_text() or "").strip()
                if text:
                    pages.append({"page": i + 1, "text": text})
        return pages
    except Exception as exc:
        logger.warning("Page-level PDF extraction failed: %s", exc)
        return []


def extract_pages_from_docx(file_bytes: bytes) -> List[Dict]:
    """
    Group every 15 paragraphs into a synthetic 'page' (DOCX has no real page numbers).
    Returns list of {"page": int, "text": str}.
    """
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        PARA_PER_PAGE = 15
        pages = []
        for i in range(0, len(paragraphs), PARA_PER_PAGE):
            text = "\n".join(paragraphs[i : i + PARA_PER_PAGE])
            pages.append({"page": (i // PARA_PER_PAGE) + 1, "text": text})
        return pages
    except Exception as exc:
        logger.warning("DOCX page extraction failed: %s", exc)
        return []


def extract_pages_from_pptx(file_bytes: bytes) -> List[Dict]:
    """Each slide = one page. Returns list of {"page": int, "text": str}."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        pages = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(r.text for r in para.runs if r.text.strip())
                        if line.strip():
                            parts.append(line.strip())
            if parts:
                pages.append({"page": i + 1, "text": "\n".join(parts)})
        return pages
    except Exception as exc:
        logger.warning("PPTX page extraction failed: %s", exc)
        return []


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_pages(
    pages: List[Dict],
    max_chars: int = _CHUNK_SIZE,
    overlap: int = _CHUNK_OVERLAP,
) -> List[Dict]:
    """
    Split each page's text into overlapping character-level chunks.
    Each chunk: {"page": int, "chapter": str, "text": str}.
    """
    chunks: List[Dict] = []
    current_chapter: Optional[str] = None

    for page_info in pages:
        page_num = page_info["page"]
        text     = page_info["text"]

        heading = _detect_heading(text)
        if heading:
            current_chapter = heading

        start = 0
        while start < len(text):
            end        = start + max_chars
            chunk_text = text[start:end].strip()
            if chunk_text:
                chunks.append({
                    "page":    page_num,
                    "chapter": current_chapter or f"Page {page_num}",
                    "text":    chunk_text,
                })
            if end >= len(text):
                break
            start = end - overlap   # overlap ensures boundary sentences aren't lost

    return chunks


# ── Embedding ─────────────────────────────────────────────────────────────────

def embed_chunks(chunks: List[Dict]):
    """
    Generate embeddings for all chunks using fastembed (BAAI/bge-small-en-v1.5).
    Returns a list of numpy vectors, or None if fastembed is unavailable.
    The caller must handle None (falls back to returning first top_k chunks).
    """
    if not chunks:
        return None
    try:
        from fastembed import TextEmbedding
        model      = TextEmbedding("BAAI/bge-small-en-v1.5")
        texts      = [c["text"] for c in chunks]
        embeddings = list(model.embed(texts))
        logger.info("RAG: embedded %d chunks", len(chunks))
        return embeddings
    except Exception as exc:
        logger.warning("fastembed unavailable (%s) — RAG will use first-k chunks", exc)
        return None


# ── Retrieval ─────────────────────────────────────────────────────────────────

def retrieve_top_chunks(
    chunks:     List[Dict],
    embeddings,          # list of numpy vectors or None
    query:      str,
    top_k:      int = 5,
) -> List[Dict]:
    """
    Retrieve the top_k most relevant chunks via cosine similarity.
    Returned chunks are sorted by page number so the AI receives coherent context.
    Falls back to the first top_k chunks when embeddings are unavailable.
    """
    if embeddings is None or not chunks:
        return chunks[:top_k]

    try:
        import numpy as np
        from fastembed import TextEmbedding

        model     = TextEmbedding("BAAI/bge-small-en-v1.5")
        q_vec     = np.array(list(model.embed([query]))[0])
        emb_mat   = np.array(embeddings)

        # Normalised cosine similarity
        row_norms = np.linalg.norm(emb_mat, axis=1, keepdims=True)
        row_norms[row_norms == 0] = 1e-10
        normed    = emb_mat / row_norms
        q_norm    = q_vec   / (np.linalg.norm(q_vec) + 1e-10)
        scores    = normed @ q_norm

        top_idx   = np.argsort(scores)[::-1][:top_k].tolist()
        # Re-sort by page order for coherent context
        top_idx   = sorted(top_idx)
        retrieved = [chunks[i] for i in top_idx]
        logger.info("RAG: retrieved %d/%d chunks for query=%r", len(retrieved), len(chunks), query[:50])
        return retrieved

    except Exception as exc:
        logger.warning("RAG retrieval error (%s) — falling back to first %d chunks", exc, top_k)
        return chunks[:top_k]


# ── Convenience: full pipeline ────────────────────────────────────────────────

def build_context_from_file(
    file_bytes: bytes,
    file_ext:   str,     # "pdf", "docx", "pptx"
    query:      str,
    top_k:      int = 6,
) -> tuple[List[Dict], str]:
    """
    End-to-end RAG pipeline for a single file.
    Returns (retrieved_chunks, context_text) where context_text is ready to
    inject into the AI prompt.  retrieved_chunks carry page/chapter metadata.
    """
    # Step 1 — extract pages
    ext = file_ext.lower().lstrip(".")
    if ext == "pdf":
        pages = extract_pages_from_pdf(file_bytes)
    elif ext == "docx":
        pages = extract_pages_from_docx(file_bytes)
    elif ext == "pptx":
        pages = extract_pages_from_pptx(file_bytes)
    else:
        pages = []

    if not pages:
        logger.warning("RAG: no pages extracted from .%s file", ext)
        return [], ""

    # Step 2 — chunk
    chunks = chunk_pages(pages)
    if not chunks:
        return [], ""

    # Step 3 — embed
    embeddings = embed_chunks(chunks)

    # Step 4 — retrieve
    top_chunks = retrieve_top_chunks(chunks, embeddings, query, top_k=top_k)

    # Step 5 — format context
    parts = []
    for c in top_chunks:
        parts.append(
            f"[Page {c['page']} | {c['chapter']}]\n{c['text']}"
        )
    context_text = "\n\n---\n\n".join(parts)

    return top_chunks, context_text
